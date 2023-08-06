from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해
from pptx.util import Cm
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from os import path

import six, copy
from ..pkg_main_module import program_execution_date
from ._utill import get_str_length
from .convert import convert_file_ppt_to_pdf


def duplicate_slide(prs, index: int):
    """
    prs 객체의 index 번째 슬라이드를 복사한 뒤에 맨 끝에 추가한다.
    - 복제한 슬라이드를 반환한다.
    """

    template = prs.slides[index]

    blank_slide_layout = prs.slide_layouts[6]

    copied_slide = prs.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(value.reltype, value._target, value.rId)

    return copied_slide


def delete_slide(prs, index: int) -> None:
    """prs 객체의 index 번째 슬라이드를 제거한다."""
    slide = prs.slides[index]
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


class PowerPoint:
    """python-pptx 라이브러리를 활용해서 파워포인트를 보조한다."""

    prs = None  # 파워포인트 객체
    file_name = None

    def __init__(self, file_name: str = None) -> None:
        """
        PowerPoint 객체를 초기화한다.
        - file_name 값은 추후에 저장할 파워포인트의 이름이다.
            - '.pptx' 로 끝나지 않으면 끝에 추가한다.
        - file_name 값의 파일이 존재하면, 해당 파일을 불러읽는다.
        - file_name 값 자체가 없으면, 'YYYYMMDD_powerpoint.pptx' 이름으로 저장된다.
        """
        # file_name 이 존재하는 파일이면 불러읽는다.
        if path.exists(file_name):
            self.prs = Presentation(file_name)
        # 존재하지 않는 파일이면, 빈 파일을 만든다.
        else:
            self.prs = Presentation()

        # 이 값은 추후에 파일을 저장할때 이름으로 사용된다.
        self.file_name = file_name

    def set_file_name(self, file_name: str) -> None:
        """
        file_name 을 변경한다.
        - '.pptx' 로 끝나지 않으면 끝에 추가한다.
        """
        if not file_name.endswith(file_name):
            file_name += None
        self.file_name = file_name

    def duplicate_slide(self, index: int):
        """
        prs 객체의 index 번째 슬라이드를 복사한 뒤에 맨 끝에 추가한다.
        - 복제한 슬라이드를 반환한다.
        """
        return duplicate_slide(self.prs, index)

    def delete_slide(self, index: int) -> None:
        """prs 객체의 index 번째 슬라이드를 제거한다."""
        delete_slide(self.prs, index)

    def save(self, *, overwrite: bool = False, convert_pdf: bool = False) -> None:
        """
        prs 객체를 파워포인트 파일로 저장한다.
        - self.file_name 값이 None이면, 'YYYYMMDD_powerpoint.pptx' 이름으로 저장된다.
        - self.file_name 값이 '.pptx' 로 끝나지 않으면 끝에 추가한다.
        - self.file_name 값이 이미 존재하는 파일 이름일때,
            - overwrite = True 이면 상관없이 저장한다.
            - overwrite = False  면 이름을 변경해서 저장하는데,
                - 앞에 YYYYMMDD_ 를 추가하고
                - 그래도 파일명이 중복되면, 끝에 '(N)' 을 추가한다. [N = 1, 2, 3, ...]
        - convert_pdf 값이 True이면, PPT파일 저장 이후, 자동으로 PDF파일로도 변환한다.
        """
        from ._os import next_path, path_join

        # - self.file_name 값이 None이면, 'YYYYMMDD_powerpoint.pptx' 이름으로 저장된다.
        if self.file_name is None:
            self.file_name = f"{program_execution_date}_powerpoint.pptx"
            src_parent_abspath = None
        else:
            file_name_splits = self.file_name.split(path.sep)
            src_parent_abspath = path.sep.join(file_name_splits[:-1])

        # - self.file_name 값이 '.pptx' 로 끝나지 않으면 끝에 추가한다.
        if not self.file_name.endswith(".pptx"):
            self.file_name += ".pptx"

        # - self.file_name 값이 이미 존재하는 파일 이름일때,
        if path.exists(self.file_name):
            # overwrite = True 이면 상관없이 저장한다.
            # overwrite = False  면 이름을 변경해서 저장하는데,
            if not overwrite:
                # 앞에 YYYYMMDD_ 를 추가하고
                file_name = file_name_splits[-1]
                src_parent_abspath = path.sep.join(file_name_splits[:-1])
                self.file_name = path_join(src_parent_abspath, f"{program_execution_date}_{file_name}")

                # 그래도 파일명이 중복되면, 끝에 '(N)' 을 추가한다. [N = 1, 2, 3, ...]
                if path.exists(self.file_name):
                    path_pattern = f"{self.file_name[:-5]} (%s).pptx"
                    self.file_name = path_join(src_parent_abspath, next_path(path_pattern))

        # prs 객체를 파워포인트 파일로 저장한다.
        self.prs.save(self.file_name)
        print(f"PowerPoint '{self.file_name}' file saved")

        if src_parent_abspath and convert_pdf:
            convert_file_ppt_to_pdf(self.file_name)
            print(f"PowerPoint '{self.file_name}' file convert to PDF file")

    def add_textbox(self, slide, text_data: str, left, top, width, height, *, length_type: str = "Cm", font_name: str = "현대하모니 L", font_pt_size: int = 16, font_rgb: list = [0x00, 0x00, 0x00], bold: bool = False):
        """
        특정 슬라이드의 지정된 위치에 텍스트 박스를 넣는다.
        - 이후, 추가작업을 위해서 생성한 textbox 를 반환한다.
        - Default 값으로,
            - length_type: length의 단위. Cm 이다.
            - 폰트 이름 : 현대하모니 L
            - 폰트 크기 : 16 [단위는 Pt]
            - 폰트 색상 : 검정 = [0x00, 0x00, 0x00]
            - 굵음 처리 : False
        """
        return add_textbox(slide, text_data, left, top, width, height, length_type=length_type, font_name=font_name, font_pt_size=font_pt_size, font_rgb=font_rgb, bold=bold)

    def add_table(self, slide, table_data, left, top, width, height=None, *, length_type: str = "Cm", header_center_align_vertical=True, header_center_align_horizontal=True, body_center_align_vertical=True, body_center_align_horizontal=False):
        """
        특정 슬라이드의 지정된 위치에 테이블을 넣는다.
        - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
        - 첫 열 역시 해더가 포함되어 있다고 가정하고, 그대로 테이블로 만든다.
        - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
        - length_type: length의 단위로, Default 값은 Cm 이다.
        - header_center_align: 표 헤더의 중앙정렬 여부로, Default 값은 True 이다.
        - header_center_align: 표 몸통의 중앙정렬 여부로, Default 값은 False 이며, 좌측정렬된다.
        """
        return add_table(
            slide,
            table_data,
            left,
            top,
            width,
            height,
            length_type=length_type,
            header_center_align_vertical=header_center_align_vertical,
            header_center_align_horizontal=header_center_align_horizontal,
            body_center_align_vertical=body_center_align_vertical,
            body_center_align_horizontal=body_center_align_horizontal,
        )

    def add_table_vertical(self, slide, table_data, left, top, width, height=None, *, length_type: str = "Cm"):
        """
        특정 슬라이드의 지정된 위치에 테이블을 넣는다.
        - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
        - 다만, add_table_vertical은 첫 열에 헤더 값이, 두번째 열에 데이터 값이 들어있는 테이블이다.
        - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
        - length_type: length의 단위로, Default 값은 Cm 이다.
        """
        return add_table_vertical(slide, table_data, left, top, width, height, length_type=length_type)

    def add_individual_tables(self, slide, table_data, left, top, width, height=None, interval=0.59, line_break_col=-1, *, length_type: str = "Cm"):
        """
        특정 슬라이드의 지정된 위치에 테이블들을 넣는다.
        - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
        - 첫 열 역시 해더가 포함되어 있다고 가정하고, 그대로 테이블로 만든다.
        - 대신 한번에 합치는 것이 아니라, interval 간격으로 테이블을 개별로 추가한다.
        - interval 은 테이블 간 간격이며, Default 값은 0.59 이다.
        - line_break_col 는 테이블에서 줄을 바꾸는 '열(column)' 값이며, Default 값은 -1 로, 맨 뒤의 열은 줄을 바꾼다.
        - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
        - length_type: length의 단위로, Default 값은 Cm 이다.
        """
        return add_individual_tables(slide, table_data, left, top, width, height, interval, line_break_col, length_type=length_type)

    def fill_table(self, table, table_data) -> None:
        """(i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다."""
        fill_table(table, table_data)

    def set_style_table(self, table, header_center_align_vertical=True, header_center_align_horizontal=True, body_center_align_vertical=True, body_center_align_horizontal=False) -> None:
        """
        주어진 table에 대해서 style을 잡는다.
        - 첫 행은 만 색깔 칠하고 굵은 글씨를 사용한다.
        - 나머지 행은 투명한 배경과 일반 글씨를 사용한다.
        - header_center_align_vertical: 표 헤더의 수직 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
        - header_center_align_horizontal: 표 헤더의 수평 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
        - body_center_align_vertical: 표 몸통의 수직 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
        - body_center_align_horizontal: 표 몸통의 수평 중앙정렬 여부로, Default 값은 False 이며, 좌측정렬된다.
        """
        set_style_table(
            table,
            header_center_align_vertical=header_center_align_vertical,
            header_center_align_horizontal=header_center_align_horizontal,
            body_center_align_vertical=body_center_align_vertical,
            body_center_align_horizontal=body_center_align_horizontal,
        )

    def auto_adjust_table_columns(self, table) -> None:
        """
        주어진 테이블 객체의 열들을 자동 조정한다.
        - 각 열에서 가장 긴 길이에 맞춰서, 열 길이를 맞춘다.
        - 열 길이는 한글은 2글자, 영문, 숫자는 1글자로 계산한다.
        """
        auto_adjust_table_columns(table)

    def set_table_column_widths(self, table, column_widths: list, length_type: str = "Cm") -> None:
        """
        주어진 테이블의 열 길이(Column_width)를 column_widths 값으로 지정한다.
        - column_widths의 값이 False 판정이면 변경하지 않는다. [None, "", 0, 등등]
        - 기본 length_type은 "Cm" 이다.
        """
        set_table_column_widths(table, column_widths, length_type)

    def set_table_column_font_size(self, table, column_idxs: list = [], font_size: int = 9, is_included_header=False) -> None:
        """
        주어진 테이블에서 column_idxs 열에 해당하는 문자들의 폰트 크기를 font_size 값으로 지정한다.
        - 기본 font_size는 9이다.
        - 기본 column_idxs는 [] 이다. 빈 값이 주어지면 모든 열에 대해서 적용한다.
        - is_included_header는 헤더(첫 행) 도 적용할지 여부로, Default 값은 False이다.
        """
        set_table_column_font_size(table, column_idxs, font_size, is_included_header)

    def merge_table_columns_common_text(self, table, column_idxs: list = [], is_blank_merge=False):
        """
        주어진 테이블에서 column_idxs 열에 해당하는 문자열 중, 연속된 문자열인 셀을 병합한다.
        - 기본 column_idxs는 [] 이다. 빈 값이 주어지면 모든 열에 대해서 적용한다.
        - 기본 is_blank_merge은 False 이다. True이면, Cell의 값이 False이면 이전 값과 동일하게 계산해서 합친다.
        """
        merge_table_columns_common_text(table, column_idxs, is_blank_merge)

    def merge_table_columns_common_text2(self, table, column_idx: int = 0, is_blank_merge=False):
        """
        주어진 테이블에서 column_idx 열부터 해당 테이블의 마지막 열 끝까지 반복한다.
        각 열에서, 첫행~마지막 행에 해당하는 문자열 중, 연속된 문자열인 셀을 병합한다.
        - 기본 column_idx는 0 이다.
        - 기본 is_blank_merge은 False 이다. True이면, Cell의 값이 False이면 이전 값과 동일하게 계산해서 합친다.
        """
        merge_table_columns_common_text2(table, column_idx, is_blank_merge)


###################################################################################
###################################################################################
# 객체 없이도 바로 사용할 수 있는 전역함수들


def add_textbox(slide, text_data: str, left, top, width, height, *, length_type: str = "Cm", font_name: str = "현대하모니 L", font_pt_size: int = 16, font_rgb: list = [0x00, 0x00, 0x00], bold: bool = False):
    """
    특정 슬라이드의 지정된 위치에 텍스트 박스를 넣는다.
    - length_type: length의 단위로, Default 값은 Cm 이다.
    - 이후, 추가작업을 위해서 생성한 textbox 를 반환한다.
    - Default 값으로,
        - 폰트 이름 : 현대하모니 L
        - 폰트 크기 : 16 [단위는 Pt]
        - 폰트 색상 : 검정 = [0x00, 0x00, 0x00]
        - 굵음 처리 : False
    """
    left = eval(f"{length_type}(left)")
    top = eval(f"{length_type}(top)")
    width = eval(f"{length_type}(width)")
    height = eval(f"{length_type}(height)")
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    tf.text = text_data

    for paragraph in tf.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_pt_size)
        paragraph.font.color.rgb = eval(f"RGBColor({str(font_rgb)[1:-1]})")
        paragraph.font.bold = bold

    return textbox


def add_table(slide, table_data, left, top, width, height=None, *, length_type: str = "Cm", header_center_align_vertical=True, header_center_align_horizontal=True, body_center_align_vertical=True, body_center_align_horizontal=False):
    """
    특정 슬라이드의 지정된 위치에 테이블을 넣는다.
    - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
    - 첫 열 역시 해더가 포함되어 있다고 가정하고, 그대로 테이블로 만든다.
    - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
    - length_type: length의 단위로, Default 값은 Cm 이다.
    - header_center_align_vertical: 표 헤더의 수직 중앙정렬 여부로, Default 값은 True 이다.
    - header_center_align_horizontal: 표 헤더의 수평 중앙정렬 여부로, Default 값은 True 이다.
    - body_center_align_vertical: 표 몸통의 수직 중앙정렬 여부로, Default 값은 False 이며, 상측정렬된다.
    - body_center_align_horizontal: 표 몸통의 수평 중앙정렬 여부로, Default 값은 False 이며, 좌측정렬된다.
    """
    rows = len(table_data)  # 행
    cols = len(table_data[0])  # 열

    if height is None:
        height = rows * 0.68

    left = eval(f"{length_type}(left)")
    top = eval(f"{length_type}(top)")
    width = eval(f"{length_type}(width)")
    height = eval(f"{length_type}(height)")

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    tbl = table_shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # NoStyleTableGrid

    # (i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다.
    fill_table(table, table_data)

    # 주어진 table에 대해서 style을 잡는다.
    set_style_table(
        table,
        header_center_align_vertical=header_center_align_vertical,
        header_center_align_horizontal=header_center_align_horizontal,
        body_center_align_vertical=body_center_align_vertical,
        body_center_align_horizontal=body_center_align_horizontal,
    )

    return table


def add_table_vertical(slide, table_data, left, top, width, height=None, *, length_type: str = "Cm"):
    """
    특정 슬라이드의 지정된 위치에 테이블을 넣는다.
    - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
    - 다만, add_table_vertical은 첫 열에 헤더 값이, 두번째 열에 데이터 값이 들어있는 테이블이다.
    - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
    - length_type: length의 단위로, Default 값은 Cm 이다.
    """
    rows = len(table_data)  # 행
    cols = 2  # 열

    if height is None:
        height = rows * 0.68

    left = eval(f"{length_type}(left)")
    top = eval(f"{length_type}(top)")
    width1 = round(width / 3, 2)
    width2 = round(width - round(width / 3, 1), 2)
    width = eval(f"{length_type}(width)")
    height = eval(f"{length_type}(height)")

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    tbl = table_shape._element.graphic.graphicData.tbl
    tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # NoStyleTableGrid

    # (i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다.
    fill_table(table, table_data)

    # 주어진 table에 대해서 style을 잡는다.
    set_style_table_vertical(table)

    # 세로 테이블은 대략 1:2의 비율로 세팅한다.
    set_table_column_widths(table, [width1, width2])

    return table


def add_individual_tables(slide, table_data, left, top, width, height=None, interval=0.59, line_break_col=-1, *, length_type: str = "Cm"):
    """
    특정 슬라이드의 지정된 위치에 테이블들을 넣는다.
    - table_data 은 n행 m열이 동일한 list[list] 값을 받는다고 가정한다.
    - 첫 열 역시 해더가 포함되어 있다고 가정하고, 그대로 테이블로 만든다.
    - 대신 한번에 합치는 것이 아니라, interval 간격으로 테이블을 개별로 추가한다.
    - interval 은 테이블 간 간격이며, Default 값은 0.59 이다.
    - line_break_col 는 테이블에서 줄을 바꾸는 '열(column)' 값이며, Default 값은 -1 로, 맨 뒤의 열은 줄을 바꾼다.
    - 이후, 추가작업을 위해서 생성한 table_shape.table 을 반환한다.
    - length_type: length의 단위로, Default 값은 Cm 이다.
    """
    tables = list()

    if line_break_col:
        headers1 = table_data[0][:line_break_col]
        headers2 = table_data[0][line_break_col:]
        len_headers2 = len(headers2)

        table_data = table_data[1:]
        table_datas1 = [table_data_[:line_break_col] for table_data_ in table_data]
        table_datas2 = [table_data_[line_break_col:] for table_data_ in table_data]

        left = eval(f"{length_type}(left)")
        width = eval(f"{length_type}(width)")

        rows = 2  # 행
        if height is None:
            height = rows * 0.68
        height = eval(f"{length_type}(height)")

        for i, table_data1 in enumerate(table_datas1):
            cols = len(headers1)  # 열
            top_ = eval(f"{length_type}(top + i * (interval + (len_headers2 + 1) * 1.36))")

            table_shape = slide.shapes.add_table(rows, cols, left, top_, width, height)
            table = table_shape.table
            tbl = table_shape._element.graphic.graphicData.tbl
            tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # NoStyleTableGrid

            # (i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다.
            fill_table(table, [headers1, table_data1])
            # 주어진 table에 대해서 style을 잡는다.
            set_style_table(table)
            tables.append(table)

            for j, header2 in enumerate(headers2, 0):
                rows = 2  # 행
                cols = 1  # 열

                top_ = eval(f"{length_type}(top + i * (interval + (len_headers2 + 1) * 1.36) + (j + 1) * 1.36)")

                table_shape = slide.shapes.add_table(rows, cols, left, top_, width, height)
                table2 = table_shape.table
                tbl = table_shape._element.graphic.graphicData.tbl
                tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # NoStyleTableGrid

                fill_table(table2, [[header2], [table_datas2[i][j]]])
                # 주어진 table에 대해서 style을 잡는다.
                set_style_table(table2)
                tables.append(table2)
    else:
        headers1 = table_data[0][:]
        headers2 = []
        len_headers2 = len(headers2)

        table_data = table_data[1:]
        table_datas1 = [table_data_[:] for table_data_ in table_data]
        table_datas2 = [[] for table_data_ in table_data]

        left = eval(f"{length_type}(left)")
        width = eval(f"{length_type}(width)")

        rows = 2  # 행
        if height is None:
            height = rows * 0.68
        height = eval(f"{length_type}(height)")

        for i, table_data1 in enumerate(table_datas1):
            cols = len(headers1)  # 열
            top_ = eval(f"{length_type}(top + i * (interval + (len_headers2 + 1) * 1.36))")

            table_shape = slide.shapes.add_table(rows, cols, left, top_, width, height)
            table = table_shape.table
            tbl = table_shape._element.graphic.graphicData.tbl
            tbl[0][-1].text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # NoStyleTableGrid

            # (i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다.
            fill_table(table, [headers1, table_data1])
            # 주어진 table에 대해서 style을 잡는다.
            set_style_table(table)
            tables.append(table)

    return tables


def fill_table(table, table_data) -> None:
    """(i 행, j열)의 크기가 동일한 table, table_data에 대해서, table_data의 값을 table로 채운다."""
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            cell_ = table.cell(row, col)
            cell_.text = table_data[row][col]


def set_style_table(table, header_center_align_vertical=True, header_center_align_horizontal=True, body_center_align_vertical=True, body_center_align_horizontal=False) -> None:
    """
    주어진 table에 대해서 style을 잡는다.
    - 첫 행은 만 색깔 칠하고 굵은 글씨를 사용한다.
    - 나머지 행은 투명한 배경과 일반 글씨를 사용한다.
    - header_center_align_vertical: 표 헤더의 수직 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
    - header_center_align_horizontal: 표 헤더의 수평 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
    - body_center_align_vertical: 표 몸통의 수직 중앙정렬 여부로, Default 값은 True 이며, 중앙정렬된다.
    - body_center_align_horizontal: 표 몸통의 수평 중앙정렬 여부로, Default 값은 False 이며, 좌측정렬된다.
    """
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            for paragraph in cell.text_frame.paragraphs:
                paragraph_font = paragraph.font
                if row_idx == 0:
                    cell.fill.fore_color.rgb = RGBColor(0xDB, 0xE4, 0xF0)
                    paragraph_font.bold = True
                    if header_center_align_vertical or header_center_align_horizontal:
                        cell_center_alignment(cell, header_center_align_vertical, header_center_align_horizontal)  # 수직, 수평 정렬

                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    if body_center_align_vertical or body_center_align_horizontal:
                        cell_center_alignment(cell, body_center_align_vertical, body_center_align_horizontal)  # 수직, 수평 정렬
                paragraph_font.fill.solid()
                # paragraph_font.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
                paragraph_font.size = Pt(10)


def cell_center_alignment(cell, center_align_vertical=True, center_align_horizontal=True):
    if center_align_vertical:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 수직 가운데 정렬
    if center_align_horizontal:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER  # 수평 가운데 정렬


def set_style_table_vertical(table) -> None:
    """
    주어진 table에 대해서 세로열 style을 잡는다.
    - 첫 열은 만 색깔 칠하고 굵은 글씨를 사용한다.
    - 둘째 열은 투명한 배경과 일반 글씨를 사용한다.
    """
    for row in table.rows:
        for col_idx, cell in enumerate(row.cells):
            cell.fill.solid()
            for paragraph in cell.text_frame.paragraphs:
                paragraph_font = paragraph.font
                if col_idx == 0:
                    cell.fill.fore_color.rgb = RGBColor(0xDB, 0xE4, 0xF0)
                    paragraph_font.bold = True
                else:
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                paragraph_font.fill.solid()
                # paragraph_font.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
                paragraph_font.size = Pt(10)


def auto_adjust_table_columns(table) -> None:
    """
    주어진 테이블 객체의 열들을 자동 조정한다.
    - 각 열에서 가장 긴 길이에 맞춰서, 열 길이를 맞춘다.
    - 열 길이는 한글은 2글자, 영문, 숫자는 1글자로 계산한다.
    """
    for col in range(len(table.columns)):
        col_text_max_length = 0
        for row in range(len(table.rows)):
            cell_ = table.cell(row, col)
            col_text_max_length = max(col_text_max_length, get_str_length(cell_.text))
        col_text_max_length = round(col_text_max_length / 4, 2)
        col_text_max_length = round(col_text_max_length, 2)
        # print("col_text_max_length:", col_text_max_length)
        table.columns[col].width = Cm(col_text_max_length)


def set_table_column_widths(table, column_widths: list, length_type: str = "Cm") -> None:
    """
    주어진 테이블의 열 길이(Column_width)를 column_widths 값으로 지정한다.
    - column_widths의 값이 False 판정이면 변경하지 않는다. [None, "", 0, 등등]
    - 기본 length_type은 "Cm" 이다.
    """
    for c, column_width in enumerate(column_widths):
        if column_width:
            table.columns[c].width = eval(f"{length_type}(column_width)")


def set_table_column_font_size(table, column_idxs: list = [], font_size: int = 9, is_included_header=False) -> None:
    """
    주어진 테이블에서 column_idxs 열에 해당하는 문자들의 폰트 크기를 font_size 값으로 지정한다.
    - 기본 font_size는 9이다.
    - 기본 column_idxs는 [] 이다. 빈 값이 주어지면 모든 열에 대해서 적용한다.
    - is_included_header는 헤더(첫 행) 도 적용할지 여부로, Default 값은 False이다.
    """
    for col in range(len(table.columns)):
        if column_idxs and col not in column_idxs:
            continue
        for row in range(len(table.rows)):
            if row == 0 and not is_included_header:
                continue
            for paragraph in table.cell(row, col).text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)


def merge_table_columns_common_text(table, column_idxs: list = [], is_blank_merge=False):
    """
    주어진 테이블에서 column_idxs 열에 해당하는 문자열 중, 연속된 문자열인 셀을 병합한다.
    - 기본 column_idxs는 [] 이다. 빈 값이 주어지면 모든 열에 대해서 적용한다.
    - 기본 is_blank_merge은 False 이다. True이면, Cell의 값이 False이면 이전 값과 동일하게 계산해서 합친다.
    """
    for col in range(len(table.columns)):
        if column_idxs and col not in column_idxs:
            continue
        first_cell = table.cell(1, col)
        first_text = first_cell.text
        prev_cell = table.cell(1, col)
        curr_text = table.cell(1, col).text
        if is_blank_merge and (not curr_text or curr_text == "None"):
            curr_text = first_text

        for row in range(2, len(table.rows)):
            curr_cell = table.cell(row, col)
            curr_text = curr_cell.text
            if is_blank_merge and (not curr_text or curr_text == "None"):
                curr_text = first_text
            if curr_text != first_text:
                # print(f"{last_text} != {curr_text}")
                if first_cell != prev_cell:
                    first_cell.merge(prev_cell)
                    first_cell.text = first_text
                    for paragraph in first_cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                    first_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                first_cell = curr_cell
                first_text = curr_cell.text
            prev_cell = curr_cell

        if curr_text == first_text:
            if first_cell != prev_cell:
                first_cell.merge(prev_cell)
                first_cell.text = first_text
                for paragraph in first_cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                first_cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def merge_table_columns_common_text2(table, column_idx: int = 0, is_blank_merge=False):
    """
    주어진 테이블에서 column_idx 열부터 해당 테이블의 마지막 열 끝까지 반복한다.
    각 열에서, 첫행~마지막 행에 해당하는 문자열 중, 연속된 문자열인 셀을 병합한다.
    - 기본 column_idx는 0 이다.
    - 기본 is_blank_merge은 False 이다. True이면, Cell의 값이 False이면 이전 값과 동일하게 계산해서 합친다.
    """
    recursive_merge_table(table, column_idx, 1, len(table.rows) - 1, is_blank_merge)


def recursive_merge_table(table, column_idx: int, start_row_idx: int, end_row_idx: int, is_blank_merge=False):
    """
    특정 열(column_idx)에 대해서, 해당 행의 주어진 행부터, 동일한 값이 반복되는 행까지의 셀을 합친다.
    """
    # 조건. 현재 열이 column_idx를 벗어나는 경우, 나간다.
    if column_idx >= len(table.columns):
        return

    # start_row_idx 행은 첫번째이자 이전 셀이다.
    first_row_idx = start_row_idx
    first_cell = table.cell(start_row_idx, column_idx)
    first_text = first_cell.text
    prev_cell = first_cell
    prev_row_idx = start_row_idx
    curr_text = ""

    # 해당 열에서, start_row_idx + 1 ~ end_row_idx 행 동안 반복한다.
    for curr_row_idx in range(start_row_idx + 1, end_row_idx + 1):
        # 현재 행의 셀을 접근한다.
        curr_cell = table.cell(curr_row_idx, column_idx)
        curr_text = curr_cell.text
        if is_blank_merge and (not curr_text or curr_text == "None"):
            curr_text = first_text

        # 현재 행의 값과, 처음 행의 값이 다르다면,
        if curr_text != first_text:
            # 처음 행과 이전행까지의 셀을 병합을 해야하는데
            # 만약 처음 행과 이전행까지의 셀이 동일한 값이라면, 병합할 필요가 없다.
            # 이 둘이 서로 다를 경우,
            if first_cell != prev_cell:

                # 1. 셀 병합을 진행한다.
                first_cell.merge(prev_cell)
                first_cell.text = first_text
                for paragraph in first_cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                first_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 2. 다음 열에서도, 해당 범위에 대해서 셀 병합을 진행한다.
                recursive_merge_table(table, column_idx + 1, first_row_idx, prev_row_idx, is_blank_merge)

            # 셀이 1개인 경우는 병합할 필요는 없고,
            # 처음 행의 셀의 글자를 재조정한다.
            first_row_idx = curr_row_idx
            first_cell = curr_cell
            first_text = curr_cell.text
        # 이전셀도 재조정한다.
        prev_cell = curr_cell
        prev_row_idx = curr_row_idx

    # 해당 열에서, 행이 모두 끝난 뒤에도, 마지막 curr_text의 값과 first_text의 값이 같았다면
    if curr_text == first_text:
        # 셀이 1개인지 여부를 정하고, 마저 병합을 진행한다.
        if first_cell != prev_cell:
            first_cell.merge(prev_cell)
            first_cell.text = first_text
            for paragraph in first_cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
            first_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            recursive_merge_table(table, column_idx + 1, first_row_idx, prev_row_idx, is_blank_merge)
